from pathlib import Path

from pptx import Presentation

from accelerator_rag.corpus.parsed_document import ParsedDocument, ParsedUnit


def _extract_slide_title(slide) -> str | None:
    title_shape = slide.shapes.title

    if title_shape is None:
        return None

    title = title_shape.text.strip()

    return title or None


def _extract_slide_text(slide) -> str:
    texts: list[str] = []

    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue

        text = shape.text.strip()

        if text:
            texts.append(text)

    return "\n".join(texts)


def _extract_slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame

    if text_frame is None:
        return ""

    return text_frame.text.strip()


def parse_pptx(path: str | Path) -> ParsedDocument:
    source_path = Path(path)

    if not source_path.exists():
        raise FileNotFoundError(source_path)

    if source_path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected .pptx file: {source_path}")

    presentation = Presentation(source_path)

    units: list[ParsedUnit] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = _extract_slide_title(slide)
        text = _extract_slide_text(slide)
        notes = _extract_slide_notes(slide)

        warnings: list[str] = []

        if not text:
            warnings.append("empty_slide_text")

        if title is None and text:
            warnings.append("missing_title_placeholder")

        units.append(
            ParsedUnit(
                index=slide_index,
                unit_type="slide",
                title=title,
                text=text,
                metadata={
                    "slide_number": slide_index,
                    "notes": notes,
                    "shape_count": len(slide.shapes),
                    "title_source": "placeholder" if title is not None else None,
                },
                warnings=warnings,
            )
        )

    return ParsedDocument(
        source_path=str(source_path),
        file_name=source_path.name,
        document_type="pptx",
        units=units,
        metadata={
            "slide_count": len(presentation.slides),
        },
    )